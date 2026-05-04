"""
Create a concise 6-slide project presentation using python-pptx.

Run:
    python scripts/create_project_presentation.py
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import List

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
RESULTS_METRICS_DIR = ROOT / "results" / "metrics"
OUTPUT_PATH = ROOT / "results" / "HAR_Project_Summary.pptx"


def _load_json(path: Path) -> dict:
    if not path.exists():
        return {}
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def _slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = title
    tf = s.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item


def _results_bullets() -> List[str]:
    pamap2_base = _load_json(RESULTS_METRICS_DIR / "pamap2_baselines.json")
    hhar_base = _load_json(RESULTS_METRICS_DIR / "HHAR_baselines.json")
    pamap2_deep = _load_json(RESULTS_METRICS_DIR / "pamap2_deep_models.json")
    hhar_deep = _load_json(RESULTS_METRICS_DIR / "hhar_deep_models.json")

    bullets = []

    if pamap2_base.get("XGBoost", {}).get("mean_accuracy") is not None:
        acc = pamap2_base["XGBoost"]["mean_accuracy"]
        bullets.append(f"PAMAP2 LOSO baseline best: XGBoost accuracy = {acc:.4f}")
    else:
        bullets.append("PAMAP2 baseline summary exists in results/metrics/pamap2_baselines.json")

    if hhar_base.get("XGBoost", {}).get("mean_accuracy") is not None:
        acc = hhar_base["XGBoost"]["mean_accuracy"]
        bullets.append(f"HHAR LOSO baseline best: XGBoost accuracy = {acc:.4f}")
    else:
        bullets.append("HHAR baseline summary exists in results/metrics/HHAR_baselines.json")

    if pamap2_deep:
        gnn = pamap2_deep.get("gnn", {}).get("accuracy")
        gnn_lstm = pamap2_deep.get("gnn_lstm", {}).get("accuracy")
        if gnn is not None and gnn_lstm is not None:
            bullets.append(
                f"PAMAP2 deep models: GNN = {gnn:.4f}, GNN+LSTM = {gnn_lstm:.4f}"
            )
        else:
            bullets.append("PAMAP2 deep model results are saved in pamap2_deep_models.json")
    else:
        bullets.append("PAMAP2 deep model metrics file is not present")

    if hhar_deep:
        available = ", ".join(sorted(hhar_deep.keys()))
        bullets.append(f"HHAR deep metrics currently recorded for: {available}")
    else:
        bullets.append("HHAR deep model metrics file is not present")

    bullets.append(
        "Additional outputs saved: holdout summaries, graph ablation, profiling, interpretability, and error analysis JSON files"
    )
    return bullets


def main() -> None:
    prs = Presentation()

    _slide(
        prs,
        "1) Problem / Objective",
        [
            "Build a continuous Human Activity Recognition (HAR) pipeline from wearable/phone sensor streams",
            "Compare classical baselines (SVM, RandomForest, XGBoost) with deep models (LSTM, GNN, GNN+LSTM)",
            "Evaluate on PAMAP2 and HHAR using subject-aware evaluation scripts in this repo",
            "Produce reproducible artifacts: metrics JSONs, predictions arrays, and saved model files",
        ],
    )

    _slide(
        prs,
        "2) Data and Preprocessing",
        [
            "Datasets: PAMAP2 (multi-IMU body sensors) and HHAR (phone/watch accelerometer streams)",
            "Data download is implemented in src/data_download.py with UCI API + ZIP fallback",
            "Preprocessing in src/preprocessing.py: resampling, per-channel normalization, sliding-window segmentation",
            "Processed outputs are saved as data/processed/{dataset}_X.npy, {dataset}_y.npy, and {dataset}_subjects.npy",
        ],
    )

    _slide(
        prs,
        "3) Approach (Baselines + LSTM)",
        [
            "Baselines in src/baselines.py extract handcrafted window features and train SVM/RF/XGBoost",
            "Deep models in src/models.py include LSTM-only, GNN-only, and GNN+LSTM variants",
            "Dataset wrappers in src/dataset.py feed flat windows, graph windows, or subject-safe sequences",
            "Graph structure and node features are built in src/graph_construction.py for sensor-dependency modeling",
        ],
    )

    _slide(
        prs,
        "4) Pipeline / Workflow",
        [
            "1) Download raw data -> 2) preprocess into .npy arrays -> 3) train/evaluate models",
            "Main LOSO deep run: scripts/run_full_pipeline.py",
            "Extended experiments (CNN1D, ablations, profiling, interpretability): scripts/experiments.py",
            "Holdout experiments and quick comparisons: scripts/train_holdout.py",
            "All run outputs are written under results/metrics and results/models",
        ],
    )

    _slide(
        prs,
        "5) Results / Current Progress",
        _results_bullets(),
    )

    _slide(
        prs,
        "6) Conclusion / Next Steps",
        [
            "Core codebase is in place: download, preprocessing, training, evaluation, experiments, and tests",
            "Current repo already stores substantial metric outputs for baseline, deep, ablation, and holdout runs",
            "Align naming/consistency across result files and scripts (e.g., HHAR filename casing, model key variants)",
            "Complete missing/placeholder scripts and regenerate final presentation plots into results/plots",
            "Use this deck as a concise checkpoint of the project's implemented state",
        ],
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
